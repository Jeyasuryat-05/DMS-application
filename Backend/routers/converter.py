"""
converter.py — Universal file-to-PDF converter.

Pure Python (no LibreOffice needed):
  Images  → Pillow → PDF
  DOCX    → python-docx + reportlab
  XLSX    → openpyxl + reportlab
  PPTX    → python-pptx + reportlab
  TXT/CSV/JSON/XML → reportlab
  DXF     → ezdxf + matplotlib → PDF

LibreOffice (optional, best quality for DWG/STEP):
  Any format → soffice --headless --convert-to pdf
"""

from fastapi import APIRouter, Depends, HTTPException, Request
from fastapi.responses import FileResponse
from sqlalchemy.orm import Session
import os, shutil, traceback

from database import get_db
import models

router    = APIRouter()
CONV_DIR  = "uploads/converted"
os.makedirs(CONV_DIR, exist_ok=True)


# ─── Token validation ──────────────────────────────────────────────────────────

def _auth(token: str, db: Session):
    try:
        from jose import jwt
        from routers.auth import SECRET_KEY, ALGORITHM
        p   = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        uid = p.get("sub")
        if uid:
            return db.query(models.User).filter(
                models.User.id == int(uid), models.User.is_active == True
            ).first()
    except Exception:
        pass
    return None


# ─── Conversion helpers ────────────────────────────────────────────────────────

def _cached(file_id, src):
    p = os.path.join(CONV_DIR, f"{file_id}.pdf")
    if os.path.exists(p) and os.path.getmtime(p) >= os.path.getmtime(src):
        return p
    return None


def _libreoffice(src, out_dir):
    """Try LibreOffice headless conversion."""
    for exe in [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/usr/bin/soffice", "/usr/bin/libreoffice",
        "soffice", "libreoffice",
    ]:
        if os.path.exists(exe) or shutil.which(exe):
            import subprocess
            try:
                subprocess.run(
                    [exe, "--headless", "--convert-to", "pdf",
                     "--outdir", out_dir, src],
                    timeout=60, capture_output=True,
                )
                base = os.path.splitext(os.path.basename(src))[0]
                pdf  = os.path.join(out_dir, base + ".pdf")
                if os.path.exists(pdf):
                    return pdf
            except Exception:
                pass
    return None


def _img_to_pdf(src, dest):
    from PIL import Image
    from reportlab.pdfgen import canvas as C
    from reportlab.lib.pagesizes import A4
    img = Image.open(src)
    if img.mode in ("RGBA", "P", "LA", "CMYK"):
        img = img.convert("RGB")
    W, H   = A4
    iw, ih = img.size
    scale  = min(W / iw, H / ih, 1.0)
    nw, nh = iw * scale, ih * scale
    c = C.Canvas(dest, pagesize=A4)
    c.drawInlineImage(img, (W-nw)/2, (H-nh)/2, nw, nh)
    c.save()


def _txt_to_pdf(src, dest):
    from reportlab.platypus import SimpleDocTemplate, Preformatted
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    with open(src, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    style = ParagraphStyle("code", fontName="Courier", fontSize=8, leading=11)
    doc   = SimpleDocTemplate(dest, pagesize=A4,
                               leftMargin=2*cm, rightMargin=2*cm,
                               topMargin=2*cm, bottomMargin=2*cm)
    doc.build([Preformatted(text, style)])


def _xlsx_to_pdf(src, dest):
    import openpyxl
    from reportlab.platypus import (SimpleDocTemplate, Table, TableStyle,
                                    Paragraph, Spacer)
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib import colors
    from reportlab.lib.units import cm

    wb     = openpyxl.load_workbook(src, data_only=True)
    styles = getSampleStyleSheet()
    story  = []
    for ws in wb.worksheets:
        story.append(Paragraph(f"Sheet: {ws.title}", styles["Heading2"]))
        story.append(Spacer(1, 6))
        rows = []
        for row in ws.iter_rows(
            max_row=min(ws.max_row or 1, 500),
            max_col=min(ws.max_column or 1, 26),
            values_only=True,
        ):
            rows.append([str(c if c is not None else "") for c in row])
        if rows:
            t = Table(rows, repeatRows=1, hAlign="LEFT")
            t.setStyle(TableStyle([
                ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#1565C0")),
                ("TEXTCOLOR",  (0,0), (-1,0), colors.white),
                ("FONTNAME",   (0,0), (-1,-1), "Helvetica"),
                ("FONTSIZE",   (0,0), (-1,-1), 7),
                ("GRID",       (0,0), (-1,-1), 0.3, colors.grey),
                ("ROWBACKGROUNDS", (0,1), (-1,-1),
                 [colors.white, colors.HexColor("#EEF4FF")]),
                ("VALIGN",     (0,0), (-1,-1), "MIDDLE"),
                ("TOPPADDING", (0,0), (-1,-1), 2),
                ("BOTTOMPADDING", (0,0), (-1,-1), 2),
            ]))
            story.append(t)
        story.append(Spacer(1, 16))
    doc = SimpleDocTemplate(dest, pagesize=landscape(A4),
                             leftMargin=1*cm, rightMargin=1*cm,
                             topMargin=1.5*cm, bottomMargin=1.5*cm)
    doc.build(story)


def _docx_to_pdf(src, dest):
    import docx
    from reportlab.platypus import (SimpleDocTemplate, Paragraph,
                                    Spacer, Table, TableStyle)
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.units import cm

    document = docx.Document(src)
    styles   = getSampleStyleSheet()
    story    = []

    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            story.append(Spacer(1, 4))
            continue
        sname = "Normal"
        pname = para.style.name if para.style else ""
        if "Heading 1" in pname:   sname = "Heading1"
        elif "Heading 2" in pname: sname = "Heading2"
        elif "Heading"   in pname: sname = "Heading3"
        try:
            story.append(Paragraph(text, styles[sname]))
        except Exception:
            story.append(Paragraph(text, styles["Normal"]))
        story.append(Spacer(1, 3))

    for table in document.tables:
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        if rows:
            t = Table(rows)
            t.setStyle(TableStyle([
                ("GRID",       (0,0), (-1,-1), 0.3, colors.grey),
                ("FONTSIZE",   (0,0), (-1,-1), 8),
                ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#1565C0")),
                ("TEXTCOLOR",  (0,0), (-1,0), colors.white),
                ("TOPPADDING", (0,0), (-1,-1), 3),
                ("BOTTOMPADDING", (0,0), (-1,-1), 3),
            ]))
            story.append(t)
            story.append(Spacer(1, 10))

    doc = SimpleDocTemplate(dest, pagesize=A4,
                             leftMargin=2*cm, rightMargin=2*cm,
                             topMargin=2*cm, bottomMargin=2*cm)
    doc.build(story)


def _pptx_to_pdf(src, dest):
    from pptx import Presentation
    from reportlab.platypus import (SimpleDocTemplate, Paragraph,
                                    Spacer, HRFlowable)
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.lib import colors
    from reportlab.lib.units import cm

    prs    = Presentation(src)
    styles = getSampleStyleSheet()
    story  = []
    for i, slide in enumerate(prs.slides, 1):
        story.append(Paragraph(f"Slide {i}", styles["Heading2"]))
        story.append(Spacer(1, 4))
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        try:
                            story.append(Paragraph(para.text, styles["Normal"]))
                        except Exception:
                            pass
                        story.append(Spacer(1, 2))
        story.append(HRFlowable(width="100%", thickness=0.5, color=colors.grey))
        story.append(Spacer(1, 12))

    doc = SimpleDocTemplate(dest, pagesize=landscape(A4),
                             leftMargin=2*cm, rightMargin=2*cm,
                             topMargin=2*cm, bottomMargin=2*cm)
    doc.build(story)


def _dxf_to_pdf(src, dest):
    import ezdxf
    from ezdxf.addons.drawing import RenderContext, Frontend
    from ezdxf.addons.drawing.matplotlib import MatplotlibBackend
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    doc = ezdxf.readfile(src)
    msp = doc.modelspace()
    fig = plt.figure(figsize=(16, 12), facecolor="white")
    ax  = fig.add_axes([0.05, 0.05, 0.9, 0.9])
    ctx = RenderContext(doc)
    out = MatplotlibBackend(ax)
    Frontend(ctx, out).draw_layout(msp, finalize=True)
    fig.savefig(dest, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _csv_to_pdf(src, dest):
    import csv
    from reportlab.platypus import (SimpleDocTemplate, Table, TableStyle, Spacer)
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.lib import colors
    from reportlab.lib.units import cm

    with open(src, "r", encoding="utf-8", errors="replace") as f:
        rows = list(csv.reader(f))

    rows = rows[:500]  # limit rows
    if not rows:
        _txt_to_pdf(src, dest)
        return

    t = Table(rows, repeatRows=1, hAlign="LEFT")
    t.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#1565C0")),
        ("TEXTCOLOR",  (0,0), (-1,0), colors.white),
        ("FONTNAME",   (0,0), (-1,-1), "Helvetica"),
        ("FONTSIZE",   (0,0), (-1,-1), 7),
        ("GRID",       (0,0), (-1,-1), 0.3, colors.grey),
        ("ROWBACKGROUNDS", (0,1), (-1,-1),
         [colors.white, colors.HexColor("#EEF4FF")]),
        ("TOPPADDING", (0,0), (-1,-1), 2),
        ("BOTTOMPADDING", (0,0), (-1,-1), 2),
    ]))
    doc = SimpleDocTemplate(dest, pagesize=landscape(A4),
                             leftMargin=1*cm, rightMargin=1*cm,
                             topMargin=1.5*cm, bottomMargin=1.5*cm)
    doc.build([t, Spacer(1, 10)])


# ─── Main conversion dispatcher ───────────────────────────────────────────────

def convert_to_pdf(file_id: int, src: str, filename: str) -> str:
    cached = _cached(file_id, src)
    if cached:
        return cached

    ext  = os.path.splitext(filename)[1].lower().lstrip(".")
    dest = os.path.join(CONV_DIR, f"{file_id}.pdf")

    # Already PDF
    if ext == "pdf":
        return src

    errors = []

    # 1. Try LibreOffice first for complex formats
    lo_exts = {"dwg","dxf","stp","step","iges","igs","doc","xls","ppt","odp","odt","ods","html","htm"}
    if ext in lo_exts:
        lo = _libreoffice(src, CONV_DIR)
        if lo:
            if lo != dest:
                shutil.move(lo, dest)
            return dest

    # 2. Pure Python converters
    try:
        if ext in ("png","jpg","jpeg","gif","bmp","webp","tiff","tif","svg"):
            _img_to_pdf(src, dest)
        elif ext in ("docx",):
            _docx_to_pdf(src, dest)
        elif ext in ("xlsx",):
            _xlsx_to_pdf(src, dest)
        elif ext in ("pptx",):
            _pptx_to_pdf(src, dest)
        elif ext in ("csv",):
            _csv_to_pdf(src, dest)
        elif ext in ("txt","log","xml","json","py","js","ts","md","css","html","htm"):
            _txt_to_pdf(src, dest)
        elif ext == "dxf":
            try:
                _dxf_to_pdf(src, dest)
            except Exception as e1:
                errors.append(f"ezdxf: {e1}")
                # DXF fallback: LibreOffice
                lo = _libreoffice(src, CONV_DIR)
                if lo:
                    if lo != dest: shutil.move(lo, dest)
                    return dest
                raise RuntimeError(f"DXF conversion failed. {'; '.join(errors)}")
        else:
            # Last resort: LibreOffice for anything else
            lo = _libreoffice(src, CONV_DIR)
            if lo:
                if lo != dest: shutil.move(lo, dest)
                return dest
            raise RuntimeError(
                f"No converter available for .{ext.upper()} files.\n\n"
                f"To enable: install LibreOffice from libreoffice.org "
                f"and restart the backend server.\n"
                f"LibreOffice supports DWG, STEP, DOC, XLS, PPT and 100+ other formats."
            )

        if not os.path.exists(dest):
            raise RuntimeError(f"Conversion produced no output file for .{ext}")

        return dest

    except RuntimeError:
        raise
    except Exception as e:
        errors.append(str(e))
        # Try LibreOffice as final fallback
        lo = _libreoffice(src, CONV_DIR)
        if lo:
            if lo != dest: shutil.move(lo, dest)
            return dest
        raise RuntimeError(
            f"Conversion failed for .{ext.upper()}.\n"
            f"Error: {'; '.join(errors)}\n\n"
            f"Fix: Run this in your Backend terminal:\n"
            f"pip install python-docx openpyxl Pillow reportlab ezdxf python-pptx\n"
            f"Then restart: python main.py"
        )


# ─── API endpoint ─────────────────────────────────────────────────────────────

@router.get("/files/{file_id}/pdf")
def get_as_pdf(
    file_id: int,
    request: Request,
    token: str = None,
    db: Session = Depends(get_db),
):
    user = _auth(token or "", db)
    if not user:
        raise HTTPException(401, "Unauthorized")

    f = db.query(models.DocumentFile).filter(models.DocumentFile.id == file_id).first()
    if not f:
        raise HTTPException(404, "File not found")
    if not os.path.exists(f.file_path):
        raise HTTPException(404, "File not found on disk")

    try:
        pdf_path = convert_to_pdf(f.id, f.file_path, f.filename or f.file_path)
    except RuntimeError as e:
        raise HTTPException(422, str(e))
    except Exception:
        raise HTTPException(500, f"Unexpected error:\n{traceback.format_exc()}")

    # Only log on the initial full request — skip PDF range/seeking requests
    if not request.headers.get("range"):
        db.add(models.AuditLog(
            document_id=f.document_id, user_id=user.id,
            action="File Viewed as PDF", note=f.filename,
        ))
        db.add(models.FileAccessLog(
            document_id=f.document_id, file_id=f.id,
            user_id=user.id, action="view",
        ))
        db.commit()

    fname = os.path.splitext(f.filename or "file")[0] + ".pdf"
    return FileResponse(
        pdf_path,
        media_type="application/pdf",
        headers={"Content-Disposition": f"inline; filename={fname}"},
    )
