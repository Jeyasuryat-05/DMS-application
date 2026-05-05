import os, shutil, traceback
from django.http import FileResponse
from rest_framework.decorators import api_view, authentication_classes, permission_classes, parser_classes
from rest_framework.parsers import JSONParser
from rest_framework.response import Response

from api.models import DocumentFile, AuditLog, FileAccessLog

CONV_DIR = 'uploads/converted'
os.makedirs(CONV_DIR, exist_ok=True)


def _cached(file_id, src):
    p = os.path.join(CONV_DIR, f'{file_id}.pdf')
    if os.path.exists(p) and os.path.getmtime(p) >= os.path.getmtime(src):
        return p
    return None


def _libreoffice(src, out_dir):
    for exe in [
        r'C:\Program Files\LibreOffice\program\soffice.exe',
        r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
        '/usr/bin/soffice', '/usr/bin/libreoffice',
        'soffice', 'libreoffice',
    ]:
        if os.path.exists(exe) or shutil.which(exe):
            import subprocess
            try:
                subprocess.run(
                    [exe, '--headless', '--convert-to', 'pdf', '--outdir', out_dir, src],
                    timeout=60, capture_output=True,
                )
                base = os.path.splitext(os.path.basename(src))[0]
                pdf = os.path.join(out_dir, base + '.pdf')
                if os.path.exists(pdf):
                    return pdf
            except Exception:
                pass
    return None


def _img_to_pdf(src, dest):
    from PIL import Image
    from reportlab.pdfgen import canvas as C
    from reportlab.lib.pagesizes import A4
    img = Image.open(src)
    if img.mode in ('RGBA', 'P', 'LA', 'CMYK'):
        img = img.convert('RGB')
    W, H = A4
    iw, ih = img.size
    scale = min(W / iw, H / ih, 1.0)
    nw, nh = iw * scale, ih * scale
    c = C.Canvas(dest, pagesize=A4)
    c.drawInlineImage(img, (W - nw) / 2, (H - nh) / 2, nw, nh)
    c.save()


def _txt_to_pdf(src, dest):
    from reportlab.platypus import SimpleDocTemplate, Preformatted
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    with open(src, 'r', encoding='utf-8', errors='replace') as f:
        text = f.read()
    style = ParagraphStyle('code', fontName='Courier', fontSize=8, leading=11)
    doc = SimpleDocTemplate(dest, pagesize=A4, leftMargin=2*cm, rightMargin=2*cm, topMargin=2*cm, bottomMargin=2*cm)
    doc.build([Preformatted(text, style)])


def _xlsx_to_pdf(src, dest):
    import openpyxl
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib import colors
    from reportlab.lib.units import cm
    wb = openpyxl.load_workbook(src, data_only=True)
    styles = getSampleStyleSheet()
    story = []
    for ws in wb.worksheets:
        story.append(Paragraph(f'Sheet: {ws.title}', styles['Heading2']))
        story.append(Spacer(1, 6))
        rows = []
        for row in ws.iter_rows(max_row=min(ws.max_row or 1, 500), max_col=min(ws.max_column or 1, 26), values_only=True):
            rows.append([str(c if c is not None else '') for c in row])
        if rows:
            t = Table(rows, repeatRows=1, hAlign='LEFT')
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1565C0')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
                ('FONTSIZE', (0, 0), (-1, -1), 7),
                ('GRID', (0, 0), (-1, -1), 0.3, colors.grey),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#EEF4FF')]),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('TOPPADDING', (0, 0), (-1, -1), 2),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 2),
            ]))
            story.append(t)
        story.append(Spacer(1, 16))
    doc = SimpleDocTemplate(dest, pagesize=landscape(A4), leftMargin=1*cm, rightMargin=1*cm, topMargin=1.5*cm, bottomMargin=1.5*cm)
    doc.build(story)


def _docx_to_pdf(src, dest):
    import docx
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.units import cm
    document = docx.Document(src)
    styles = getSampleStyleSheet()
    story = []
    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            story.append(Spacer(1, 4))
            continue
        sname = 'Normal'
        pname = para.style.name if para.style else ''
        if 'Heading 1' in pname:
            sname = 'Heading1'
        elif 'Heading 2' in pname:
            sname = 'Heading2'
        elif 'Heading' in pname:
            sname = 'Heading3'
        try:
            story.append(Paragraph(text, styles[sname]))
        except Exception:
            story.append(Paragraph(text, styles['Normal']))
        story.append(Spacer(1, 3))
    for table in document.tables:
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        if rows:
            t = Table(rows)
            t.setStyle(TableStyle([
                ('GRID', (0, 0), (-1, -1), 0.3, colors.grey),
                ('FONTSIZE', (0, 0), (-1, -1), 8),
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1565C0')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ]))
            story.append(t)
            story.append(Spacer(1, 10))
    doc = SimpleDocTemplate(dest, pagesize=A4, leftMargin=2*cm, rightMargin=2*cm, topMargin=2*cm, bottomMargin=2*cm)
    doc.build(story)


def _pptx_to_pdf(src, dest):
    from pptx import Presentation
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.lib import colors
    from reportlab.lib.units import cm
    prs = Presentation(src)
    styles = getSampleStyleSheet()
    story = []
    for i, slide in enumerate(prs.slides, 1):
        story.append(Paragraph(f'Slide {i}', styles['Heading2']))
        story.append(Spacer(1, 4))
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        try:
                            story.append(Paragraph(para.text, styles['Normal']))
                        except Exception:
                            pass
                        story.append(Spacer(1, 2))
        story.append(HRFlowable(width='100%', thickness=0.5, color=colors.grey))
        story.append(Spacer(1, 12))
    doc = SimpleDocTemplate(dest, pagesize=landscape(A4), leftMargin=2*cm, rightMargin=2*cm, topMargin=2*cm, bottomMargin=2*cm)
    doc.build(story)


def _csv_to_pdf(src, dest):
    import csv
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Spacer
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.lib import colors
    from reportlab.lib.units import cm
    with open(src, 'r', encoding='utf-8', errors='replace') as f:
        rows = list(csv.reader(f))
    rows = rows[:500]
    if not rows:
        _txt_to_pdf(src, dest)
        return
    t = Table(rows, repeatRows=1, hAlign='LEFT')
    t.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1565C0')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
        ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
        ('FONTSIZE', (0, 0), (-1, -1), 7),
        ('GRID', (0, 0), (-1, -1), 0.3, colors.grey),
        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#EEF4FF')]),
        ('TOPPADDING', (0, 0), (-1, -1), 2),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 2),
    ]))
    doc = SimpleDocTemplate(dest, pagesize=landscape(A4), leftMargin=1*cm, rightMargin=1*cm, topMargin=1.5*cm, bottomMargin=1.5*cm)
    doc.build([t, Spacer(1, 10)])


def convert_to_pdf(file_id, src, filename):
    cached = _cached(file_id, src)
    if cached:
        return cached

    ext = os.path.splitext(filename)[1].lower().lstrip('.')
    dest = os.path.join(CONV_DIR, f'{file_id}.pdf')

    if ext == 'pdf':
        return src

    errors = []
    lo_exts = {'dwg', 'dxf', 'stp', 'step', 'iges', 'igs', 'doc', 'xls', 'ppt', 'odp', 'odt', 'ods', 'html', 'htm'}
    if ext in lo_exts:
        lo = _libreoffice(src, CONV_DIR)
        if lo:
            if lo != dest:
                shutil.move(lo, dest)
            return dest

    try:
        if ext in ('png', 'jpg', 'jpeg', 'gif', 'bmp', 'webp', 'tiff', 'tif', 'svg'):
            _img_to_pdf(src, dest)
        elif ext == 'docx':
            _docx_to_pdf(src, dest)
        elif ext == 'xlsx':
            _xlsx_to_pdf(src, dest)
        elif ext == 'pptx':
            _pptx_to_pdf(src, dest)
        elif ext == 'csv':
            _csv_to_pdf(src, dest)
        elif ext in ('txt', 'log', 'xml', 'json', 'py', 'js', 'ts', 'md', 'css', 'html', 'htm'):
            _txt_to_pdf(src, dest)
        else:
            lo = _libreoffice(src, CONV_DIR)
            if lo:
                if lo != dest:
                    shutil.move(lo, dest)
                return dest
            raise RuntimeError(
                f'No converter available for .{ext.upper()} files.\n\n'
                'To enable: install LibreOffice from libreoffice.org and restart the backend server.'
            )

        if not os.path.exists(dest):
            raise RuntimeError(f'Conversion produced no output file for .{ext}')
        return dest

    except RuntimeError:
        raise
    except Exception as e:
        errors.append(str(e))
        lo = _libreoffice(src, CONV_DIR)
        if lo:
            if lo != dest:
                shutil.move(lo, dest)
            return dest
        raise RuntimeError(
            f'Conversion failed for .{ext.upper()}.\n'
            f'Error: {"; ".join(errors)}\n\n'
            'Fix: pip install python-docx openpyxl Pillow reportlab ezdxf python-pptx'
        )


@api_view(['GET'])
def get_as_pdf(request, file_id):
    try:
        f = DocumentFile.objects.get(id=file_id)
    except DocumentFile.DoesNotExist:
        return Response({'error': 'File not found'}, status=404)
    if not os.path.exists(f.file_path):
        return Response({'error': 'File not found on disk'}, status=404)

    try:
        pdf_path = convert_to_pdf(f.id, f.file_path, f.filename or f.file_path)
    except RuntimeError as e:
        return Response({'error': str(e)}, status=422)
    except Exception:
        return Response({'error': f'Unexpected error:\n{traceback.format_exc()}'}, status=500)

    if not request.META.get('HTTP_RANGE'):
        AuditLog.objects.create(
            document_id=f.document_id, user_id=request.user.id,
            action='File Viewed as PDF', note=f.filename,
        )
        FileAccessLog.objects.create(
            document_id=f.document_id, file_id=f.id,
            user_id=request.user.id, action='view',
        )

    fname = os.path.splitext(f.filename or 'file')[0] + '.pdf'
    response = FileResponse(open(pdf_path, 'rb'), content_type='application/pdf')
    response['Content-Disposition'] = f'inline; filename="{fname}"'
    return response
